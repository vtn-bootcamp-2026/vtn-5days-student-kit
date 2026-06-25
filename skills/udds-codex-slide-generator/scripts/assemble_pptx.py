import json
import argparse
import os
import sys
import shutil
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


SLIDE_IMAGE_PATTERN = re.compile(r"^slide(?P<number>\d+)(?P<suffix>[A-Za-z]*)_redesigned\.(?P<ext>jpe?g|png)$", re.IGNORECASE)


class UDDSAssembler:
    def __init__(self):
        pass

    def slide_sort_key(self, path):
        match = SLIDE_IMAGE_PATTERN.match(path.name)
        if not match:
            return (sys.maxsize, 1, path.name.lower())
        suffix = match.group("suffix").upper()
        return (int(match.group("number")), 0 if suffix == "" else 1, suffix, path.name.lower())

    def load_archetypes(self, review_file):
        archetypes = {}
        if not review_file:
            return archetypes

        review_path = Path(review_file)
        if not review_path.exists():
            print(f"[!] Review file not found; assembling by folder order only: {review_path}")
            return archetypes

        with open(review_path, "r", encoding="utf-8") as f:
            for line in f:
                if not line.startswith("| page_"):
                    continue
                parts = [p.strip() for p in line.split("|") if p.strip()]
                if len(parts) < 2:
                    continue
                page_token = parts[0].replace("page_", "").replace(".png", "")
                archetypes[page_token.lower()] = parts[1]
        return archetypes

    def list_slide_images(self, redesigned_folder):
        folder = Path(redesigned_folder)
        images = [
            path
            for path in folder.iterdir()
            if path.is_file() and SLIDE_IMAGE_PATTERN.match(path.name)
        ]
        return sorted(images, key=self.slide_sort_key)

    def handle_record_json(self, redesigned_folder, mode):
        folder = Path(redesigned_folder)
        records = sorted(folder.glob("slide*_redesigned.*.record.json"))
        if not records:
            return

        if mode == "keep":
            print(f"[*] Keeping {len(records)} slide record JSON file(s) in place.")
            return

        if mode == "archive":
            records_dir = folder / "records"
            records_dir.mkdir(parents=True, exist_ok=True)
            for record in records:
                shutil.move(str(record), str(records_dir / record.name))
            print(f"[*] Archived {len(records)} slide record JSON file(s) to: {records_dir}")
            return

        if mode == "delete":
            for record in records:
                record.unlink()
            print(f"[*] Deleted {len(records)} slide record JSON file(s).")
            return

        raise ValueError(f"Unsupported record JSON mode: {mode}")

    def build(self, review_file, redesigned_folder, output_pptx, record_json_mode="archive"):
        prs = Presentation()
        
        # 16:9 Layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        archetypes = self.load_archetypes(review_file)
        slide_images = self.list_slide_images(redesigned_folder)

        if not slide_images:
            print(f"[!] Error: Cannot assemble PPTX. No slide*_redesigned.jpg/png images found in: {redesigned_folder}")
            sys.exit(1)

        print(f"[*] Building PPTX with {len(slide_images)} redesigned slide image(s) from folder order...")

        for i, img_path in enumerate(slide_images):
            # Add slide (Layout 6 is blank)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

            match = SLIDE_IMAGE_PATTERN.match(img_path.name)
            slide_token = f"{match.group('number')}{match.group('suffix')}".lower() if match else img_path.stem
            archetype = archetypes.get(slide_token, "folder-order")
            print(f"  [+] Slide {i+1}: Inserted {img_path.name} ({archetype})")

        prs.save(output_pptx)
        print(f"[OK] Premium Redesigned PPTX created: {output_pptx}")
        self.handle_record_json(redesigned_folder, record_json_mode)

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--review", help="Optional path to archetype_review.md. Used only for archetype labels; slide order comes from --img.")
    parser.add_argument("--img", required=True, help="Folder containing redesigned images")
    parser.add_argument("--output", required=True, help="Path to save output PPTX")
    parser.add_argument(
        "--record-json-mode",
        choices=["archive", "delete", "keep"],
        default="archive",
        help="What to do with slide*.record.json files after assembling the PPTX. Defaults to archive.",
    )
    
    args = parser.parse_args()
    assembler = UDDSAssembler()
    assembler.build(args.review, args.img, args.output, args.record_json_mode)
