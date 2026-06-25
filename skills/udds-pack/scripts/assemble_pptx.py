import json
import argparse
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

class UDDSAssembler:
    def __init__(self):
        pass

    def build(self, review_file, redesigned_folder, output_pptx):
        prs = Presentation()
        
        # 16:9 Layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Load mappings from review file
        mappings = []
        with open(review_file, "r", encoding="utf-8") as f:
            for line in f:
                if line.startswith("| page_"):
                    parts = [p.strip() for p in line.split("|") if p.strip()]
                    mappings.append({"page": parts[0], "archetype": parts[1]})

        print(f"[*] Verifying {len(mappings)} redesigned slides...")
        
        missing_slides = []
        for item in mappings:
            page_num = item["page"].replace("page_", "").replace(".png", "")
            img_filename = f"slide{page_num}_redesigned.jpg"
            img_path = Path(redesigned_folder) / img_filename
            if not img_path.exists():
                missing_slides.append(img_filename)
                
        if missing_slides:
            print(f"[!] Warning: Missing {len(missing_slides)} redesigned images. They will be skipped in this run.")

        print(f"[*] Building PPTX with available redesigned slides...")

        for item in mappings:
            page_num = item["page"].replace("page_", "").replace(".png", "")
            img_filename = f"slide{page_num}_redesigned.jpg"
            img_path = Path(redesigned_folder) / img_filename
            
            if not img_path.exists():
                continue
                
            # Add slide (Layout 6 is blank)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
            print(f"  [+] Slide {page_num}: Inserted Redesigned Image ({item['archetype']})")

        prs.save(output_pptx)
        print(f"[OK] Premium Redesigned PPTX created: {output_pptx}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--review", required=True, help="Path to archetype_review.md")
    parser.add_argument("--img", required=True, help="Folder containing redesigned images")
    parser.add_argument("--output", required=True, help="Path to save output PPTX")
    
    args = parser.parse_args()
    assembler = UDDSAssembler()
    assembler.build(args.review, args.img, args.output)